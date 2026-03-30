"""Helper script for python-pptx operations that AppleScript can't do.
Called by the MCP server via execFile (no shell injection risk).

Usage:
  python3 pptx-helper.py set-levels <file> <slide_num> <shape_idx> <json_paragraphs>

  json_paragraphs format: [{"text": "Line 1", "level": 0}, {"text": "Sub item", "level": 1}, ...]
"""
from __future__ import annotations

import json
import sys

from pptx import Presentation


def set_levels(file: str, slide_num: int, shape_idx: int, paragraphs: list[dict]) -> None:
    prs = Presentation(file)
    slide = prs.slides[slide_num - 1]
    shape = slide.shapes[shape_idx - 1]
    tf = shape.text_frame
    tf.clear()

    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para["text"]
        p.level = para.get("level", 0)

    prs.save(file)
    print(json.dumps({"ok": True, "paragraphs": len(paragraphs)}))


def add_image(file: str, slide_num: int, image_path: str, x: int, y: int, width: int, height: int) -> None:
    from pptx.util import Pt

    prs = Presentation(file)
    slide = prs.slides[slide_num - 1]
    slide.shapes.add_picture(image_path, Pt(x), Pt(y), Pt(width), Pt(height))
    prs.save(file)
    print(json.dumps({"ok": True, "slide": slide_num, "image": image_path}))


if __name__ == "__main__":
    cmd = sys.argv[1]
    if cmd == "set-levels":
        file = sys.argv[2]
        slide_num = int(sys.argv[3])
        shape_idx = int(sys.argv[4])
        paragraphs = json.loads(sys.argv[5])
        set_levels(file, slide_num, shape_idx, paragraphs)
    elif cmd == "add-image":
        add_image(sys.argv[2], int(sys.argv[3]), sys.argv[4], int(sys.argv[5]), int(sys.argv[6]), int(sys.argv[7]), int(sys.argv[8]))
    else:
        print(json.dumps({"error": f"Unknown command: {cmd}"}))
        sys.exit(1)
